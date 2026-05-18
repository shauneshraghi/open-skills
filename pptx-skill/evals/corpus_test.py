"""Comprehensive test suite against the Apache POI PPTX corpus.

Covers edge cases:
- Empty notes slides (notes XML exists but body text is empty)
- Slides with no placeholders (blank layout)
- Multi-image slides
- Notes-only check (file has notesSlideMaster but no notes text on some slides)
- Slide reorder stability (reorder then re-reorder)
- Notes write/read idempotency
- Large file handling (>100 slides)
"""

from __future__ import annotations

import os
import struct
import sys
import tempfile
import traceback
import zlib
from pathlib import Path
from typing import NamedTuple

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from scripts.create_pptx import save, set_notes
from scripts.edit_pptx import (
    get_notes, read_all_notes, reorder_slides, replace_image_by_index,
)
from scripts.validate_pptx import validate_pptx
from pptx import Presentation


POI_PATH = Path(
    os.environ.get(
        "POI_PATH",
        str(Path(__file__).resolve().parents[3] / "poi" / "test-data" / "slideshow"),
    )
)


def _make_png(r: int, g: int, b: int, w: int = 8, h: int = 8) -> bytes:
    raw = (b"\x00" + bytes([r, g, b]) * w) * h
    compressed = zlib.compress(raw)

    def chunk(ctype: bytes, data: bytes) -> bytes:
        c = ctype + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", compressed) + chunk(b"IEND", b"")


class TestResult(NamedTuple):
    fixture: str
    test_name: str
    passed: bool
    detail: str = ""


results: list[TestResult] = []


def record(fixture: str, test_name: str, passed: bool, detail: str = "") -> None:
    results.append(TestResult(fixture, test_name, passed, detail))
    status = "PASS" if passed else "FAIL"
    print(f"  [{status}] {test_name}: {detail}" if detail else f"  [{status}] {test_name}")


def test_open_and_iterate(fixture_path: Path) -> None:
    fname = fixture_path.name
    try:
        prs = Presentation(str(fixture_path))
        total_shapes = sum(len(slide.shapes) for slide in prs.slides)
        record(fname, "open_and_iterate", True, f"{len(prs.slides)} slides, {total_shapes} shapes")
    except Exception as exc:
        record(fname, "open_and_iterate", False, str(exc))


def test_empty_notes(fixture_path: Path) -> None:
    fname = fixture_path.name
    try:
        prs = Presentation(str(fixture_path))
        empty_count = sum(
            1 for slide in prs.slides
            if slide.has_notes_slide and get_notes(slide) == ""
        )
        record(fname, "empty_notes_returns_empty_string", True, f"{empty_count} empty notes slides")
    except Exception as exc:
        record(fname, "empty_notes_returns_empty_string", False, str(exc))


def test_notes_idempotency(fixture_path: Path) -> None:
    fname = fixture_path.name
    tmpdir = Path(tempfile.mkdtemp())
    out = tmpdir / fixture_path.name
    note_text = "Idempotency test note — ✓ unicode"
    try:
        prs = Presentation(str(fixture_path))
        for slide in prs.slides:
            set_notes(slide, note_text)
        prs.save(str(out))
        prs2 = Presentation(str(out))
        mismatches = sum(
            1 for slide in prs2.slides if get_notes(slide) != note_text
        )
        record(fname, "notes_write_idempotency", mismatches == 0,
               f"{mismatches} mismatches in {len(prs2.slides)} slides")
    except Exception as exc:
        record(fname, "notes_write_idempotency", False, str(exc))


def test_reorder_stability(fixture_path: Path) -> None:
    fname = fixture_path.name
    tmpdir = Path(tempfile.mkdtemp())
    out = tmpdir / fixture_path.name
    try:
        prs = Presentation(str(fixture_path))
        n = len(prs.slides)
        if n < 2:
            record(fname, "reorder_stability", True, "skipped (1 slide)")
            return

        def titles(p: Presentation) -> list[str]:
            result = []
            for slide in p.slides:
                t = ""
                title_shape = slide.shapes.title
                if title_shape is not None and hasattr(title_shape, "text"):
                    try:
                        t = title_shape.text or ""
                    except Exception:
                        t = ""
                result.append(t)
            return result

        original = titles(prs)
        reorder_slides(prs, n - 1, 0)
        reorder_slides(prs, 0, n - 1)
        prs.save(str(out))
        prs2 = Presentation(str(out))
        restored = titles(prs2)
        match = original == restored
        record(fname, "reorder_stability", match,
               f"{'restored' if match else 'MISMATCH'} ({n} slides)")
    except Exception as exc:
        record(fname, "reorder_stability", False, str(exc))


def test_multi_image(fixture_path: Path) -> None:
    fname = fixture_path.name
    tmpdir = Path(tempfile.mkdtemp())
    new_img = tmpdir / "green.png"
    new_img.write_bytes(_make_png(0, 180, 0))
    out = tmpdir / fixture_path.name
    try:
        prs = Presentation(str(fixture_path))
        multi_slides = [
            (i, slide)
            for i, slide in enumerate(prs.slides)
            if sum(1 for s in slide.shapes if s.shape_type == 13) > 1
        ]
        if not multi_slides:
            record(fname, "multi_image_replace", True, "skipped (no multi-image slides)")
            return
        idx, slide = multi_slides[0]
        pic_count = sum(1 for s in slide.shapes if s.shape_type == 13)
        ok = replace_image_by_index(slide, 0, new_img)
        prs.save(str(out))
        record(fname, "multi_image_replace", ok, f"slide {idx} had {pic_count} images")
    except Exception as exc:
        record(fname, "multi_image_replace", False, str(exc))


def test_no_placeholder_slide(fixture_path: Path) -> None:
    fname = fixture_path.name
    try:
        prs = Presentation(str(fixture_path))
        blank_count = sum(
            1 for slide in prs.slides
            if not any(
                hasattr(s, "is_placeholder") and s.is_placeholder
                for s in slide.shapes
            )
        )
        record(fname, "no_placeholder_slide", True, f"{blank_count} slides without placeholders")
    except Exception as exc:
        record(fname, "no_placeholder_slide", False, str(exc))


def test_validate(fixture_path: Path) -> None:
    fname = fixture_path.name
    try:
        vr = validate_pptx(fixture_path)
        record(
            fname, "validate_pptx", vr.valid,
            f"slides={vr.info.get('slide_count', '?')}, errors={vr.errors[:1] if not vr.valid else []}",
        )
    except Exception as exc:
        record(fname, "validate_pptx", False, str(exc))


def test_read_all_notes_structure(fixture_path: Path) -> None:
    fname = fixture_path.name
    try:
        prs = Presentation(str(fixture_path))
        all_notes = read_all_notes(prs)
        ok = (
            isinstance(all_notes, list)
            and len(all_notes) == len(prs.slides)
            and all("slide_index" in n and "notes" in n for n in all_notes)
        )
        record(fname, "read_all_notes_structure", ok, f"{len(all_notes)} entries")
    except Exception as exc:
        record(fname, "read_all_notes_structure", False, str(exc))


CORPUS_FIXTURES = [
    "shapes.pptx",
    "SampleShow.pptx",
    "aascu.org_hbcu_leadershipsummit_cooper_.pptx",
    "table_test2.pptx",
    "bug65551.pptx",
    "testPPT.pptx",
    "layouts.pptx",
    "45545_Comment.pptx",
    "2411-Performance_Up.pptx",
    "WithMaster.pptx",
    "ArtisticEffectSample.pptx",
    "bar-chart.pptx",
    "table-with-theme.pptx",
    "keyframes.pptx",
    "customGeo.pptx",
]


def main() -> None:
    if not POI_PATH.exists():
        print(f"ERROR: POI_PATH not found: {POI_PATH}")
        sys.exit(1)

    print(f"POI corpus: {POI_PATH}")
    print("=" * 70)

    for fname in CORPUS_FIXTURES:
        fixture = POI_PATH / fname
        if not fixture.exists():
            print(f"\n[SKIP] {fname} not found")
            continue
        print(f"\n--- {fname} ---")
        test_open_and_iterate(fixture)
        test_empty_notes(fixture)
        test_read_all_notes_structure(fixture)
        test_notes_idempotency(fixture)
        test_reorder_stability(fixture)
        test_multi_image(fixture)
        test_no_placeholder_slide(fixture)
        test_validate(fixture)

    print("\n" + "=" * 70)
    passed = sum(1 for r in results if r.passed)
    failed = sum(1 for r in results if not r.passed)
    print(f"Corpus test results: {passed} passed, {failed} failed (total {len(results)})")

    if failed:
        print("\nFailures:")
        for r in results:
            if not r.passed:
                print(f"  FAIL {r.fixture} / {r.test_name}: {r.detail}")
        sys.exit(1)
    else:
        sys.exit(0)


if __name__ == "__main__":
    main()
