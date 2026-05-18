"""Benchmark pptx-skill operations against the Apache POI test corpus.

Design:
- Uses POI_PATH env var (default: /home/user/poi/test-data/slideshow) per the
  cross-platform requirement; all paths via pathlib.Path.
- tempfile.mkdtemp() for all output to avoid hardcoded paths.
- Measures wall-clock time for creation, reading, notes I/O, and round-trip edits.
"""

from __future__ import annotations

import argparse
import os
import struct
import tempfile
import time
import zlib
from pathlib import Path
from typing import NamedTuple

from pptx import Presentation
from pptx.util import Inches, Pt


def _make_png(r: int = 200, g: int = 100, b: int = 50, w: int = 10, h: int = 10) -> bytes:
    raw = (b"\x00" + bytes([r, g, b]) * w) * h
    compressed = zlib.compress(raw)

    def chunk(ctype: bytes, data: bytes) -> bytes:
        c = ctype + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", compressed) + chunk(b"IEND", b"")


class BenchResult(NamedTuple):
    name: str
    iterations: int
    total_seconds: float
    unit: str = "file"

    @property
    def per_unit(self) -> float:
        return self.total_seconds / max(self.iterations, 1)

    def __str__(self) -> str:
        return (
            f"{self.name:<45} {self.iterations:>4} {self.unit}(s)  "
            f"{self.total_seconds:6.3f}s total  "
            f"{self.per_unit * 1000:7.2f}ms/{self.unit}"
        )


def bench_creation(n: int = 50) -> BenchResult:
    from scripts.create_pptx import (
        add_content_slide, add_title_slide, create_presentation, save, set_notes,
    )
    tmpdir = Path(tempfile.mkdtemp())
    start = time.perf_counter()
    for i in range(n):
        prs = create_presentation()
        add_title_slide(prs, f"Deck {i}", "Subtitle text")
        s = add_content_slide(prs, "Content", "Body text here.")
        set_notes(s, "Speaker note for the content slide.")
        save(prs, tmpdir / f"bench_{i}.pptx")
    return BenchResult("create (title+content+notes)", n, time.perf_counter() - start)


def bench_bulk_creation(slides_per_deck: int = 20, n: int = 10) -> BenchResult:
    from scripts.create_pptx import bulk_add_slides, create_presentation, save
    tmpdir = Path(tempfile.mkdtemp())
    specs = [
        {"type": "title" if i == 0 else "content", "title": f"Slide {i}", "content": f"Content {i}"}
        for i in range(slides_per_deck)
    ]
    start = time.perf_counter()
    for i in range(n):
        prs = create_presentation()
        bulk_add_slides(prs, specs)
        save(prs, tmpdir / f"bulk_{i}.pptx")
    return BenchResult(f"bulk create ({slides_per_deck} slides/deck)", n, time.perf_counter() - start)


def bench_read_corpus(poi_path: Path) -> BenchResult:
    files = list(poi_path.glob("*.pptx"))
    ok = 0
    start = time.perf_counter()
    for f in files:
        try:
            prs = Presentation(str(f))
            for slide in prs.slides:
                _ = [s.name for s in slide.shapes]
            ok += 1
        except Exception:
            pass
    return BenchResult("read POI corpus (iterate slides)", ok, time.perf_counter() - start)


def bench_notes_write(poi_path: Path) -> BenchResult:
    from scripts.edit_pptx import set_notes
    tmpdir = Path(tempfile.mkdtemp())
    files = [f for f in poi_path.glob("*.pptx") if f.stat().st_size < 2 * 1024 * 1024][:20]
    ok = 0
    start = time.perf_counter()
    for f in files:
        try:
            prs = Presentation(str(f))
            for i, slide in enumerate(prs.slides):
                set_notes(slide, f"Benchmark note for slide {i}.")
            prs.save(str(tmpdir / f.name))
            ok += 1
        except Exception:
            pass
    return BenchResult("notes write (POI files)", ok, time.perf_counter() - start)


def bench_notes_read(poi_path: Path) -> BenchResult:
    from scripts.edit_pptx import read_all_notes
    files = [f for f in poi_path.glob("*.pptx") if f.stat().st_size < 2 * 1024 * 1024][:30]
    ok = 0
    start = time.perf_counter()
    for f in files:
        try:
            prs = Presentation(str(f))
            _ = read_all_notes(prs)
            ok += 1
        except Exception:
            pass
    return BenchResult("notes read (POI files)", ok, time.perf_counter() - start)


def bench_image_round_trip(n: int = 20) -> BenchResult:
    from scripts.create_pptx import add_blank_slide, add_picture_to_slide, create_presentation, save
    from scripts.edit_pptx import replace_image_by_index
    tmpdir = Path(tempfile.mkdtemp())
    img1 = tmpdir / "img1.png"
    img2 = tmpdir / "img2.png"
    img1.write_bytes(_make_png(255, 0, 0))
    img2.write_bytes(_make_png(0, 0, 255))
    start = time.perf_counter()
    for i in range(n):
        prs = create_presentation()
        s = add_blank_slide(prs)
        add_picture_to_slide(s, img1, alt_text="Red image")
        out1 = tmpdir / f"rt_{i}_v1.pptx"
        save(prs, out1)
        prs2 = Presentation(str(out1))
        replace_image_by_index(prs2.slides[0], 0, img2)
        save(prs2, tmpdir / f"rt_{i}_v2.pptx")
    return BenchResult("image round-trip (add+replace)", n, time.perf_counter() - start)


def bench_reorder(n: int = 30) -> BenchResult:
    from scripts.create_pptx import add_content_slide, add_title_slide, create_presentation, save
    from scripts.edit_pptx import reorder_slides
    tmpdir = Path(tempfile.mkdtemp())
    prs = create_presentation()
    for i in range(10):
        if i == 0:
            add_title_slide(prs, f"Slide {i}")
        else:
            add_content_slide(prs, f"Slide {i}")
    base = tmpdir / "reorder_base.pptx"
    save(prs, base)
    start = time.perf_counter()
    for i in range(n):
        prs2 = Presentation(str(base))
        reorder_slides(prs2, 9, 0)
        save(prs2, tmpdir / f"reorder_{i}.pptx")
    return BenchResult("slide reorder (10-slide deck)", n, time.perf_counter() - start)


def main() -> None:
    default_poi = os.environ.get(
        "POI_PATH",
        str(Path(__file__).resolve().parents[3] / "poi" / "test-data" / "slideshow"),
    )
    parser = argparse.ArgumentParser(description="Benchmark pptx-skill operations")
    parser.add_argument("--poi-path", type=Path, default=Path(default_poi))
    parser.add_argument("--quick", action="store_true", help="Run shorter iterations")
    args = parser.parse_args()

    n = 10 if args.quick else 50
    print(f"{'Benchmark':<45} {'Iters':>5}  {'Total':>8}  {'Per-unit':>12}")
    print("-" * 80)

    results = []
    if args.poi_path.exists():
        results += [
            bench_read_corpus(args.poi_path),
            bench_notes_read(args.poi_path),
            bench_notes_write(args.poi_path),
        ]
    else:
        print(f"[SKIP] POI corpus not found at {args.poi_path}")

    results += [
        bench_creation(n),
        bench_bulk_creation(20, max(n // 5, 5)),
        bench_image_round_trip(max(n // 5, 5)),
        bench_reorder(n),
    ]

    for r in results:
        print(r)


if __name__ == "__main__":
    import sys
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
    main()
