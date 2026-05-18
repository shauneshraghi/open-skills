"""Run the 5 evaluation test cases for the pptx-creation-editing skill.

Uses tempfile.mkdtemp() for all output and POI_PATH env var for corpus fixtures.
"""

from __future__ import annotations

import os
import struct
import sys
import tempfile
import traceback
import zlib
from pathlib import Path
from typing import Callable

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from scripts.create_pptx import (
    add_blank_slide, add_content_slide, add_picture_to_slide, add_title_slide,
    create_presentation, save, set_notes,
)
from scripts.edit_pptx import (
    get_alt_text, get_notes, read_all_notes, replace_image_by_index, reorder_slides,
)
from scripts.validate_pptx import validate_pptx
from pptx import Presentation


POI_PATH = Path(
    os.environ.get(
        "POI_PATH",
        str(Path(__file__).resolve().parents[3] / "poi" / "test-data" / "slideshow"),
    )
)


def _make_png(r: int, g: int, b: int, w: int = 20, h: int = 20) -> bytes:
    raw = (b"\x00" + bytes([r, g, b]) * w) * h
    compressed = zlib.compress(raw)

    def chunk(ctype: bytes, data: bytes) -> bytes:
        c = ctype + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", compressed) + chunk(b"IEND", b"")


class EvalResult:
    def __init__(self, test_id: str):
        self.test_id = test_id
        self.passed = True
        self.failures: list[str] = []
        self.log: list[str] = []

    def assert_eq(self, label: str, actual, expected) -> None:
        if actual != expected:
            self.passed = False
            self.failures.append(f"FAIL [{label}]: expected {expected!r}, got {actual!r}")
        else:
            self.log.append(f"PASS [{label}]")

    def assert_true(self, label: str, value: bool) -> None:
        if not value:
            self.passed = False
            self.failures.append(f"FAIL [{label}]: expected True, got False")
        else:
            self.log.append(f"PASS [{label}]")

    def assert_contains(self, label: str, haystack: str, needle: str) -> None:
        if needle not in haystack:
            self.passed = False
            self.failures.append(f"FAIL [{label}]: {needle!r} not in {haystack!r}")
        else:
            self.log.append(f"PASS [{label}]")


def test_create_title_and_image_slide() -> EvalResult:
    r = EvalResult("create_title_and_image_slide")
    tmpdir = Path(tempfile.mkdtemp())
    out = tmpdir / "test1.pptx"
    img_path = tmpdir / "orange.png"
    img_path.write_bytes(_make_png(255, 140, 0))

    prs = create_presentation()
    add_title_slide(prs, "Quarterly Review", "Q4 2025")
    s2 = add_blank_slide(prs)
    add_picture_to_slide(s2, img_path, left=1.0, top=1.0, width=4.0, height=3.0, alt_text="A vibrant orange square")
    set_notes(s2, "This slide contains an embedded image.")
    save(prs, out)

    prs2 = Presentation(str(out))
    r.assert_eq("slide count", len(prs2.slides), 2)
    r.assert_eq("title text", prs2.slides[0].shapes.title.text, "Quarterly Review")
    pics = [s for s in prs2.slides[1].shapes if s.shape_type == 13]
    r.assert_true("slide 1 has a picture", len(pics) > 0)
    if pics:
        r.assert_eq("alt text", get_alt_text(pics[0]), "A vibrant orange square")
    r.assert_eq("notes text", get_notes(prs2.slides[1]), "This slide contains an embedded image.")
    r.assert_true("validate_pptx passes", validate_pptx(out).valid)
    return r


def test_add_speaker_notes() -> EvalResult:
    r = EvalResult("add_speaker_notes")
    tmpdir = Path(tempfile.mkdtemp())
    out = tmpdir / "test2.pptx"
    note_text = "Welcome to the presentation. Please hold questions until the end."

    prs = create_presentation()
    add_title_slide(prs, "Opening Slide")
    add_content_slide(prs, "Details", "Some body text.")
    save(prs, out)

    prs2 = Presentation(str(out))
    r.assert_true("has_notes_slide before", not prs2.slides[0].has_notes_slide)
    set_notes(prs2.slides[0], note_text)
    r.assert_true("has_notes_slide after set", prs2.slides[0].has_notes_slide)
    prs2.save(str(out))

    prs3 = Presentation(str(out))
    r.assert_eq("notes text round-trip", get_notes(prs3.slides[0]), note_text)
    r.assert_eq("slide 1 notes unaffected", get_notes(prs3.slides[1]), "")
    return r


def test_read_notes_from_poi_fixture() -> EvalResult:
    r = EvalResult("read_notes_from_poi_fixture")
    fixture = POI_PATH / "SampleShow.pptx"
    if not fixture.exists():
        r.passed = False
        r.failures.append(f"POI fixture not found: {fixture}")
        return r

    prs = Presentation(str(fixture))
    r.assert_eq("slide count", len(prs.slides), 2)
    all_notes = read_all_notes(prs)
    r.assert_eq("read_all_notes length", len(all_notes), 2)
    r.assert_eq("slide 0 notes", all_notes[0]["notes"], "I am the notes of the first slide")
    r.assert_contains("slide 1 notes", all_notes[1]["notes"], "These are the notes of the 2nd slide")
    return r


def test_reorder_slides() -> EvalResult:
    r = EvalResult("reorder_slides")
    tmpdir = Path(tempfile.mkdtemp())
    out = tmpdir / "test4.pptx"

    prs = create_presentation()
    from pptx.util import Inches
    for label in ("Slide ONE", "Slide TWO", "Slide THREE"):
        s = add_blank_slide(prs)
        txb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txb.text_frame.text = label
    save(prs, out)

    prs2 = Presentation(str(out))
    r.assert_eq("slide count before", len(prs2.slides), 3)
    r.assert_eq("order before [0]", prs2.slides[0].shapes[0].text_frame.text, "Slide ONE")
    reorder_slides(prs2, 2, 0)
    prs2.save(str(out))

    prs3 = Presentation(str(out))
    r.assert_eq("slide count after", len(prs3.slides), 3)
    r.assert_eq("new order [0]", prs3.slides[0].shapes[0].text_frame.text, "Slide THREE")
    r.assert_eq("new order [1]", prs3.slides[1].shapes[0].text_frame.text, "Slide ONE")
    r.assert_eq("new order [2]", prs3.slides[2].shapes[0].text_frame.text, "Slide TWO")
    r.assert_true("validate_pptx passes", validate_pptx(out).valid)
    return r


def test_round_trip_image_replace_and_notes() -> EvalResult:
    r = EvalResult("round_trip_image_replace_and_notes")
    fixture = POI_PATH / "shapes.pptx"
    if not fixture.exists():
        r.passed = False
        r.failures.append(f"POI fixture not found: {fixture}")
        return r

    tmpdir = Path(tempfile.mkdtemp())
    out = tmpdir / "test5_roundtrip.pptx"
    blue_png = tmpdir / "blue.png"
    blue_png.write_bytes(_make_png(0, 0, 255))

    prs = Presentation(str(fixture))
    original_slide_count = len(prs.slides)
    slide0 = prs.slides[0]
    pics = [s for s in slide0.shapes if s.shape_type == 13]
    replaced = False
    if pics:
        replaced = replace_image_by_index(slide0, 0, blue_png)
    note_text = "Round-trip test complete."
    set_notes(slide0, note_text)
    prs.save(str(out))

    prs2 = Presentation(str(out))
    r.assert_eq("slide count unchanged", len(prs2.slides), original_slide_count)
    if pics:
        r.assert_true("image was replaced", replaced)
    r.assert_eq("notes text after reload", get_notes(prs2.slides[0]), note_text)
    r.assert_true("slide 0 has_notes_slide", prs2.slides[0].has_notes_slide)

    import zipfile
    with zipfile.ZipFile(out) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        r.assert_true("media/ has at least one image", len(media) > 0)

    r.assert_true("validate_pptx passes", validate_pptx(out).valid)
    return r


TESTS: list[Callable[[], EvalResult]] = [
    test_create_title_and_image_slide,
    test_add_speaker_notes,
    test_read_notes_from_poi_fixture,
    test_reorder_slides,
    test_round_trip_image_replace_and_notes,
]


def run_all() -> list[EvalResult]:
    results = []
    for test_fn in TESTS:
        print(f"\n{'='*60}")
        print(f"Running: {test_fn.__name__}")
        print("-" * 60)
        try:
            result = test_fn()
        except Exception as exc:
            result = EvalResult(test_fn.__name__)
            result.passed = False
            result.failures.append(f"EXCEPTION: {exc}")
            traceback.print_exc()

        for line in result.log:
            print(f"  {line}")
        for line in result.failures:
            print(f"  {line}")

        status = "PASSED" if result.passed else "FAILED"
        print(f"\n  → {status}: {result.test_id}")
        results.append(result)

    print(f"\n{'='*60}")
    passed = sum(1 for r in results if r.passed)
    print(f"Results: {passed}/{len(results)} tests passed")
    return results


if __name__ == "__main__":
    results = run_all()
    sys.exit(0 if all(r.passed for r in results) else 1)
