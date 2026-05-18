"""Create PowerPoint presentations from scratch.

Design decisions from primary sources:
- Presentation() with no args → default 10×7.5-in blank deck (python-pptx docs §Presentations)
- slide_layouts[0..6] are the standard Office theme layouts; index 0 = Title Slide,
  1 = Title and Content, 6 = Blank (python-pptx docs §Working with Slides)
- Alt text lives on <p:cNvPr descr="…"> (ECMA-376 §20.1.2.2.8); python-pptx exposes
  no high-level API for it, so we set it via lxml directly.
- Tables are wrapped in a <p:graphicFrame> (ECMA-376 §19.3.1.21); python-pptx's
  add_table() returns the GraphicFrame, and .table gives the Table object.
"""

from __future__ import annotations

import argparse
import json
import os
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ───────────────────────────────── helpers ────────────────────────────────────

def _set_alt_text(shape, alt_text: str) -> None:
    """Set accessibility alt text on any shape.

    ECMA-376 §20.1.2.2.8: the `descr` attribute on <p:cNvPr> / <a:cNvPr>
    carries the description used by assistive technologies.  python-pptx 1.0.x
    has no dedicated setter, so we write it through the element tree directly.
    """
    for attr in ("nvPicPr", "nvSpPr", "nvGraphicFramePr"):
        nvPr = getattr(shape._element, attr, None)
        if nvPr is not None:
            cNvPr = getattr(nvPr, "cNvPr", None)
            if cNvPr is not None:
                cNvPr.set("descr", alt_text)
                return


# ───────────────────────────────── creation ────────────────────────────────────

def create_presentation(
    width_inches: float = 13.33,
    height_inches: float = 7.5,
) -> Presentation:
    """Return a new blank Presentation with the given slide dimensions.

    Default 13.33×7.5 in matches the widescreen 16:9 preset (cx=12192000 EMU,
    cy=6858000 EMU).  ECMA-376 §19.2.1.23 specifies sldSz as the presentation
    element attribute governing slide canvas dimensions.
    """
    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)
    return prs


def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
) -> Any:
    """Add a title slide (layout 0: 'Title Slide') and return the Slide."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle:
        try:
            slide.placeholders[1].text = subtitle
        except (KeyError, IndexError):
            pass
    return slide


def add_content_slide(
    prs: Presentation,
    title: str,
    content: str = "",
) -> Any:
    """Add a 'Title and Content' slide (layout 1) and return the Slide."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if content:
        try:
            slide.placeholders[1].text = content
        except (KeyError, IndexError):
            pass
    return slide


def add_blank_slide(prs: Presentation) -> Any:
    """Add a blank slide (layout 6) and return the Slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_picture_to_slide(
    slide: Any,
    image_path: str | Path,
    left: float = 1.0,
    top: float = 1.5,
    width: float = 4.0,
    height: float = 3.0,
    alt_text: str = "",
) -> Any:
    """Add a picture and optionally set its accessibility alt text.

    Alt text path: <p:pic><p:nvPicPr><p:cNvPr descr="…"> (ECMA-376 §20.1.2.2.8).
    python-pptx creates the blipFill relationship automatically; the embedded
    image is referenced via `r:embed` on <a:blip> pointing to the image part.
    """
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    if alt_text:
        _set_alt_text(pic, alt_text)
    return pic


def add_text_box(
    slide: Any,
    text: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 4.0,
    height: float = 1.0,
    font_size: int = 18,
    bold: bool = False,
    color: tuple[int, int, int] | None = None,
    align: str = "left",
) -> Any:
    """Add a free-floating text box to a slide.

    python-pptx docs §Adding a Text Box: add_textbox() returns a TextBox shape;
    text formatting is via TextFrame → Paragraph → Run → Font.
    """
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    para.alignment = align_map.get(align.lower(), PP_ALIGN.LEFT)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txb


def add_table(
    slide: Any,
    rows: int,
    cols: int,
    data: list[list[str]] | None = None,
    left: float = 0.5,
    top: float = 1.5,
    width: float = 9.0,
    height: float = 2.0,
) -> Any:
    """Add a table and optionally populate it with data.

    python-pptx docs §Working with Tables: add_table() returns a GraphicFrame;
    .table gives the Table object.  ECMA-376 §20.1.4 specifies <a:tbl> structure.
    """
    gf = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = gf.table
    if data:
        for r_idx, row_data in enumerate(data[:rows]):
            for c_idx, val in enumerate(row_data[:cols]):
                tbl.cell(r_idx, c_idx).text = str(val)
    return gf


def set_notes(slide: Any, text: str) -> None:
    """Set the speaker notes text on a slide.

    python-pptx docs §Working with Notes: slide.notes_slide returns the
    NotesSlidePart, creating it (and a notesMaster if missing) on first access.
    The notes text body placeholder has ph type='body' idx=1 (ECMA-376 §19.3.1.43).

    For direct lxml manipulation (e.g. formatting, empty-paragraph insertion)::

        ns_elm = slide.notes_slide._element
        from pptx.oxml.ns import qn
        body_sp = next(
            sp for sp in ns_elm.spTree.iter(qn('p:sp'))
            if sp.nvSpPr.nvPr.get_or_add_ph().type == 'body'
        )
        txBody = body_sp.txBody
        # manipulate txBody paragraphs directly via lxml
    """
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def bulk_add_slides(prs: Presentation, slides_data: list[dict]) -> list:
    """Add multiple slides from a list of spec dicts and return the slide list.

    Each dict supports keys:
        type (str): 'title' | 'content' | 'blank'  (default: 'content')
        title (str): slide title
        content (str): body text
        image_path (str | Path): path to image file
        alt_text (str): accessibility description for the image
        notes (str): speaker notes text
        table_data (list[list[str]]): 2-D list for a table
        table_rows (int), table_cols (int): table dimensions
    """
    added = []
    for spec in slides_data:
        slide_type = spec.get("type", "content")
        title = spec.get("title", "")
        content = spec.get("content", "")

        if slide_type == "title":
            slide = add_title_slide(prs, title, content)
        elif slide_type == "blank":
            slide = add_blank_slide(prs)
        else:
            slide = add_content_slide(prs, title, content)

        if "image_path" in spec:
            add_picture_to_slide(
                slide,
                spec["image_path"],
                alt_text=spec.get("alt_text", ""),
            )

        if "table_data" in spec:
            tdata = spec["table_data"]
            rows = spec.get("table_rows", len(tdata))
            cols = spec.get("table_cols", max(len(r) for r in tdata) if tdata else 1)
            add_table(slide, rows, cols, tdata)

        if "notes" in spec:
            set_notes(slide, spec["notes"])

        added.append(slide)
    return added


def save(prs: Presentation, path: str | Path) -> Path:
    """Save the presentation; creates parent directories as needed."""
    out = Path(path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


# ───────────────────────────────── CLI ───────────────────────────────────────────

def _build_demo(out_path: Path, image_path: Path | None = None) -> None:
    prs = create_presentation()
    add_title_slide(prs, "Demo Presentation", "Created with pptx-creation-editing")
    s2 = add_content_slide(prs, "Content Slide", "Bullet point text goes here.")
    set_notes(s2, "These are the speaker notes for slide 2.")
    if image_path and image_path.exists():
        s3 = add_blank_slide(prs)
        add_picture_to_slide(s3, image_path, alt_text="Sample image")
        set_notes(s3, "Slide with an embedded image.")
    add_table(
        prs.slides[1] if len(prs.slides) > 1 else add_blank_slide(prs),
        3, 3,
        [["Name", "Score", "Grade"], ["Alice", "95", "A"], ["Bob", "82", "B"]],
        left=1.0, top=3.5, width=8.0, height=2.0,
    )
    save(prs, out_path)
    print(f"Saved: {out_path}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Create a demo .pptx file")
    parser.add_argument("--out", type=Path, default=Path("demo.pptx"), help="Output path")
    parser.add_argument("--image", type=Path, default=None, help="Optional image to embed")
    parser.add_argument(
        "--spec",
        type=Path,
        default=None,
        help="JSON file with slides_data list (bulk creation)",
    )
    args = parser.parse_args()

    if args.spec:
        with open(args.spec) as f:
            spec = json.load(f)
        prs = create_presentation()
        bulk_add_slides(prs, spec)
        save(prs, args.out)
        print(f"Saved {len(prs.slides)} slides to {args.out}")
    else:
        _build_demo(args.out, args.image)


if __name__ == "__main__":
    main()
