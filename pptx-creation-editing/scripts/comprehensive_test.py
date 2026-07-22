"""Full Apache POI corpus workflow for pptx-creation-editing.

This script automatically:
- clones a sparse, shallow Apache POI working copy into a temp directory
- uses Apache POI test-data/slideshow as the baseline corpus
- verifies required fixtures exist and are intact
- runs the PPTX regression suite against that baseline
- emits verbose structured logs and metadata
- cleans up the clone and temp outputs unless artifacts are retained

Exact upstream baseline path within apache/poi:
- test-data/slideshow/
"""

from __future__ import annotations

import argparse
import base64
import datetime as dt
import hashlib
import json
import shutil
import struct
import subprocess
import sys
import tempfile
import time
import traceback
import zipfile
import zlib
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).parent))
import create_pptx
import edit_pptx
import validate_pptx
import vba_runner

TEMP_ROOT = Path(tempfile.gettempdir()) / "kilo"
DEFAULT_REPO_URL = "https://github.com/apache/poi.git"
DEFAULT_REF = "trunk"
BASELINE_RELATIVE_PATH = Path("test-data/slideshow")
REQUIRED_FIXTURES = (
    "shapes.pptx",
    "SampleShow.pptx",
    "aascu.org_hbcu_leadershipsummit_cooper_.pptx",
    "table_test2.pptx",
    "bug65551.pptx",
    "testPPT.pptx",
    "layouts.pptx",
    "45545_Comment.pptx",
    "2411-Performance_Up.pptx",
    "WithMaster.pptx",
    "ArtisticEffectSample.pptx",
    "bar-chart.pptx",
    "table-with-theme.pptx",
    "keyframes.pptx",
    "customGeo.pptx",
)
DEFAULT_TIMEOUT_SECONDS = 240
DEFAULT_RETRIES = 3
PASS = "PASS"
FAIL = "FAIL"
SKIP = "SKIP"
JPEG_1X1 = base64.b64decode(
    "/9j/4AAQSkZJRgABAQAAAQABAAD/2wCEAAkGBxAQEBAQEA8QEA8QDw8PDw8PDw8PDw8PFREWFhURFRUYHSggGBolGxUVITEhJSkrLi4uFx8zODMsNygtLisBCgoKDg0OGxAQGy0lICYtLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLS0tLf/AABEIAAEAAgMBIgACEQEDEQH/xAAXAAEBAQEAAAAAAAAAAAAAAAAAAQID/8QAFBABAAAAAAAAAAAAAAAAAAAAAP/aAAwDAQACEAMQAAAB6AAAAP/EABQQAQAAAAAAAAAAAAAAAAAAACD/2gAIAQEAAT8Af//EABQRAQAAAAAAAAAAAAAAAAAAACD/2gAIAQIBAT8Af//EABQRAQAAAAAAAAAAAAAAAAAAACD/2gAIAQMBAT8Af//Z"
)
PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADUlEQVR42mP8z/C/HwAFgwJ/l6s5WQAAAABJRU5ErkJggg=="
)
GIF_1X1 = base64.b64decode("R0lGODdhAQABAIEAAP8AAAAAAAAAAAAAACwAAAAAAQABAAAIBAABBAQAOw==")


class WorkflowError(RuntimeError):
    """Raised when a workflow step fails in a controlled way."""


class TestSkipped(RuntimeError):
    """Raised when a test has no eligible baseline fixture to exercise."""


@dataclass
class TestRecord:
    name: str
    status: str
    detail: str
    started_at: str
    finished_at: str
    duration_seconds: float


@dataclass
class RunContext:
    args: argparse.Namespace
    run_root: Path
    repo_dir: Path
    output_dir: Path
    log_path: Path
    metadata_path: Path
    repo_url: str
    ref: str
    clone_started_at: str | None = None
    clone_finished_at: str | None = None
    repo_commit: str | None = None
    baseline_root: Path | None = None
    fixture_metadata: list[dict[str, Any]] = field(default_factory=list)
    test_records: list[TestRecord] = field(default_factory=list)
    step_logs: list[dict[str, Any]] = field(default_factory=list)
    cleanup_completed: bool = False
    test_command: str = ""


def now_utc() -> dt.datetime:
    return dt.datetime.now(dt.timezone.utc)


def iso_now() -> str:
    return now_utc().isoformat()


def log_event(ctx: RunContext, level: str, event: str, **fields: Any) -> None:
    record = {
        "timestamp": iso_now(),
        "level": level,
        "event": event,
        **fields,
    }
    ctx.step_logs.append(record)
    line = json.dumps(record, sort_keys=True, default=str)
    print(line)
    ctx.log_path.parent.mkdir(parents=True, exist_ok=True)
    with ctx.log_path.open("a", encoding="utf-8") as fh:
        fh.write(line + "\n")


def _json_safe(value: Any) -> Any:
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, dt.datetime):
        return value.isoformat()
    if isinstance(value, TestRecord):
        return value.__dict__
    if isinstance(value, dict):
        return {str(k): _json_safe(v) for k, v in value.items()}
    if isinstance(value, (list, tuple)):
        return [_json_safe(v) for v in value]
    return value


def write_metadata(ctx: RunContext, status: str, error: str | None = None) -> None:
    passed = sum(1 for rec in ctx.test_records if rec.status == PASS)
    failed = sum(1 for rec in ctx.test_records if rec.status == FAIL)
    skipped = sum(1 for rec in ctx.test_records if rec.status == SKIP)
    metadata = {
        "baseline": {
            "repo_url": ctx.repo_url,
            "ref": ctx.ref,
            "commit": ctx.repo_commit,
            "path": str(BASELINE_RELATIVE_PATH).replace("\\", "/"),
            "clone_started_at": ctx.clone_started_at,
            "clone_finished_at": ctx.clone_finished_at,
            "fixtures": ctx.fixture_metadata,
        },
        "run": {
            "command": ctx.test_command,
            "started_at": ctx.step_logs[0]["timestamp"] if ctx.step_logs else None,
            "finished_at": iso_now(),
            "status": status,
            "error": error,
            "summary": {
                "total": len(ctx.test_records),
                "passed": passed,
                "failed": failed,
                "skipped": skipped,
            },
            "tests": [_json_safe(rec) for rec in ctx.test_records],
            "logs": _json_safe(ctx.step_logs),
        },
        "cleanup": {
            "completed": ctx.cleanup_completed,
        },
    }
    ctx.metadata_path.parent.mkdir(parents=True, exist_ok=True)
    with ctx.metadata_path.open("w", encoding="utf-8") as fh:
        json.dump(metadata, fh, indent=2)


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def run_command(
    ctx: RunContext,
    command: list[str],
    *,
    cwd: Path | None = None,
    timeout: int | None = None,
    retries: int | None = None,
    retry_delay: float = 2.0,
) -> subprocess.CompletedProcess[str]:
    timeout = timeout or ctx.args.timeout
    retries = retries or ctx.args.retries
    attempt = 0
    last_error: str | None = None
    while attempt < retries:
        attempt += 1
        started = time.monotonic()
        log_event(
            ctx,
            "info",
            "command.start",
            attempt=attempt,
            retries=retries,
            cwd=str(cwd) if cwd else None,
            timeout_seconds=timeout,
            command=command,
        )
        try:
            completed = subprocess.run(
                command,
                cwd=str(cwd) if cwd else None,
                check=False,
                capture_output=True,
                text=True,
                timeout=timeout,
            )
        except subprocess.TimeoutExpired as exc:
            last_error = f"Timed out after {timeout}s: {' '.join(command)}"
            log_event(
                ctx,
                "error",
                "command.timeout",
                attempt=attempt,
                duration_seconds=round(time.monotonic() - started, 3),
                stdout=exc.stdout,
                stderr=exc.stderr,
                error=last_error,
            )
        except OSError as exc:
            last_error = f"Failed to launch command {' '.join(command)}: {exc}"
            log_event(
                ctx,
                "error",
                "command.oserror",
                attempt=attempt,
                duration_seconds=round(time.monotonic() - started, 3),
                error=last_error,
            )
        else:
            duration = round(time.monotonic() - started, 3)
            log_event(
                ctx,
                "info" if completed.returncode == 0 else "warning",
                "command.finish",
                attempt=attempt,
                duration_seconds=duration,
                returncode=completed.returncode,
                stdout=completed.stdout,
                stderr=completed.stderr,
            )
            if completed.returncode == 0:
                return completed
            last_error = (
                f"Command failed with exit code {completed.returncode}: {' '.join(command)}\n"
                f"STDERR:\n{completed.stderr.strip()}\nSTDOUT:\n{completed.stdout.strip()}"
            )
        if attempt < retries:
            sleep_for = retry_delay * attempt
            log_event(ctx, "warning", "command.retry", attempt=attempt, sleep_seconds=sleep_for, error=last_error)
            time.sleep(sleep_for)
    raise WorkflowError(last_error or f"Command failed: {' '.join(command)}")


def _make_png(r: int = 0, g: int = 128, b: int = 255, w: int = 8, h: int = 8) -> bytes:
    def chunk(tag: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + bytes([r, g, b]) * w for _ in range(h))
    idat = zlib.compress(raw)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")


def _make_bmp(r: int = 0, g: int = 128, b: int = 255, w: int = 2, h: int = 2) -> bytes:
    row_bytes = ((w * 3 + 3) // 4) * 4
    image_size = row_bytes * h
    file_size = 54 + image_size
    header = struct.pack("<2sIHHI", b"BM", file_size, 0, 0, 54)
    dib = struct.pack("<IIIHHIIIIII", 40, w, h, 1, 24, 0, image_size, 2835, 2835, 0, 0)
    pixel = bytes([b, g, r])
    row = pixel * w + (b"\x00" * (row_bytes - (w * 3)))
    return header + dib + (row * h)


def _fixture(ctx: RunContext, relative_path: str) -> Path:
    assert ctx.baseline_root is not None
    return ctx.baseline_root / relative_path


def _output_path(ctx: RunContext, fixture_name: str, name: str) -> Path:
    fixture_dir = ctx.output_dir / Path(fixture_name).stem
    fixture_dir.mkdir(parents=True, exist_ok=True)
    return fixture_dir / name


def _write_binary(path: Path, data: bytes) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(data)
    return path


def _slide_titles(prs: Presentation) -> list[str]:
    titles: list[str] = []
    for slide in prs.slides:
        title = ""
        title_shape = slide.shapes.title
        if title_shape is not None and hasattr(title_shape, "text"):
            title = title_shape.text or ""
        if not title:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text or ""
                    if text:
                        title = text
                        break
        titles.append(title)
    return titles


def _picture_shapes(slide: Any) -> list[Any]:
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _find_multi_image_slide(prs: Presentation) -> tuple[int, Any, list[Any]]:
    for index, slide in enumerate(prs.slides):
        pictures = _picture_shapes(slide)
        if len(pictures) > 1:
            return index, slide, pictures
    raise TestSkipped("no slide with more than one picture")


def _replacement_bytes_for_ext(ext: str) -> tuple[str, bytes]:
    normalized = ext.lower().lstrip(".")
    if normalized in {"jpg", "jpeg", "jpe"}:
        return ".jpg", JPEG_1X1
    if normalized == "png":
        return ".png", PNG_1X1
    if normalized == "gif":
        return ".gif", GIF_1X1
    if normalized == "bmp":
        return ".bmp", _make_bmp()
    raise TestSkipped(f"unsupported picture format for safe replacement: {normalized}")


def prepare_poi_baseline(ctx: RunContext) -> None:
    ctx.clone_started_at = iso_now()
    run_command(
        ctx,
        [
            "git",
            "clone",
            "--depth",
            "1",
            "--filter=blob:none",
            "--sparse",
            ctx.repo_url,
            str(ctx.repo_dir),
        ],
    )
    run_command(
        ctx,
        [
            "git",
            "-C",
            str(ctx.repo_dir),
            "sparse-checkout",
            "set",
            str(BASELINE_RELATIVE_PATH).replace("\\", "/"),
        ],
    )
    run_command(ctx, ["git", "-C", str(ctx.repo_dir), "checkout", ctx.ref])
    ctx.repo_commit = run_command(ctx, ["git", "-C", str(ctx.repo_dir), "rev-parse", "HEAD"], timeout=60).stdout.strip()
    ctx.clone_finished_at = iso_now()
    ctx.baseline_root = ctx.repo_dir / BASELINE_RELATIVE_PATH
    log_event(
        ctx,
        "info",
        "baseline.ready",
        repo_dir=str(ctx.repo_dir),
        baseline_root=str(ctx.baseline_root),
        commit=ctx.repo_commit,
    )


def verify_baseline(ctx: RunContext) -> None:
    assert ctx.baseline_root is not None
    if not ctx.baseline_root.exists() or not ctx.baseline_root.is_dir():
        raise WorkflowError(f"Missing baseline directory: {ctx.baseline_root}")

    fixture_metadata: list[dict[str, Any]] = []
    for fixture_name in REQUIRED_FIXTURES:
        path = _fixture(ctx, fixture_name)
        if not path.exists():
            raise WorkflowError(f"Missing required baseline fixture: {path}")
        if path.stat().st_size <= 0:
            raise WorkflowError(f"Baseline fixture is empty: {path}")
        try:
            with zipfile.ZipFile(path) as zf:
                bad = zf.testzip()
                if bad is not None:
                    raise WorkflowError(f"Fixture ZIP integrity failed for {path}: first bad entry {bad}")
        except zipfile.BadZipFile as exc:
            raise WorkflowError(f"Fixture is not a valid ZIP package: {path} ({exc})") from exc

        validation = validate_pptx.validate_pptx(path)
        if not validation.valid:
            raise WorkflowError(f"Fixture failed validate_pptx: {path}: {validation.errors}")

        fixture_metadata.append(
            {
                "path": fixture_name,
                "size": path.stat().st_size,
                "sha256": sha256_file(path),
                "zip_valid": True,
                "validation_valid": validation.valid,
                "slide_count": validation.info.get("slide_count"),
                "notes_slide_count": validation.info.get("notes_slide_count"),
                "image_count": validation.info.get("image_count"),
            }
        )
    ctx.fixture_metadata = fixture_metadata
    log_event(ctx, "info", "baseline.verified", fixture_count=len(fixture_metadata), fixtures=fixture_metadata)


class TestSuite:
    def __init__(self, ctx: RunContext) -> None:
        self.ctx = ctx

    def run_test(self, name: str, fn: Callable[[], str]) -> None:
        started = time.monotonic()
        started_at = iso_now()
        log_event(self.ctx, "info", "test.start", name=name)
        try:
            detail = fn() or ""
            status = PASS
        except TestSkipped as exc:
            detail = str(exc)
            status = SKIP
        except AssertionError as exc:
            detail = str(exc)
            status = FAIL
        except Exception as exc:
            detail = f"{type(exc).__name__}: {exc}\n{traceback.format_exc(limit=5)}"
            status = FAIL

        finished_at = iso_now()
        record = TestRecord(
            name=name,
            status=status,
            detail=detail,
            started_at=started_at,
            finished_at=finished_at,
            duration_seconds=round(time.monotonic() - started, 3),
        )
        self.ctx.test_records.append(record)
        event = "test.pass" if status == PASS else "test.skip" if status == SKIP else "test.fail"
        level = "info" if status in {PASS, SKIP} else "error"
        log_event(self.ctx, level, event, name=name, detail=detail, status=status)

    def run_fixture_test(self, fixture_name: str, label: str, fn: Callable[[Path], str]) -> None:
        path = _fixture(self.ctx, fixture_name)
        self.run_test(f"{fixture_name}: {label}", lambda: fn(path))

    def run(self) -> None:
        for fixture_name in REQUIRED_FIXTURES:
            self.run_fixture_test(fixture_name, "open_and_iterate", self.test_open_and_iterate)
            self.run_fixture_test(fixture_name, "empty_notes_returns_empty_string", self.test_empty_notes)
            self.run_fixture_test(fixture_name, "read_all_notes_structure", self.test_read_all_notes_structure)
            self.run_fixture_test(fixture_name, "notes_write_idempotency", self.test_notes_idempotency)
            self.run_fixture_test(fixture_name, "reorder_stability", self.test_reorder_stability)
            self.run_fixture_test(fixture_name, "multi_image_replace", self.test_multi_image)
            self.run_fixture_test(fixture_name, "no_placeholder_slide", self.test_no_placeholder_slide)
            self.run_fixture_test(fixture_name, "validate_pptx", self.test_validate)
        self.run_test("G1: check_powerpnt returns a string", self.test_g1)
        self.run_test("G2: generate_vbs_runner produces valid VBScript", self.test_g2)
        self.run_test("G3: write_bas_file produces valid .bas module", self.test_g3)

    def test_open_and_iterate(self, fixture_path: Path) -> str:
        prs = Presentation(str(fixture_path))
        total_shapes = sum(len(slide.shapes) for slide in prs.slides)
        return f"slides={len(prs.slides)} shapes={total_shapes}"

    def test_empty_notes(self, fixture_path: Path) -> str:
        prs = Presentation(str(fixture_path))
        empty_notes = 0
        for slide in prs.slides:
            notes = edit_pptx.get_notes(slide)
            assert isinstance(notes, str), "get_notes did not return a string"
            if slide.has_notes_slide and notes == "":
                empty_notes += 1
        return f"empty_notes_slides={empty_notes} total_slides={len(prs.slides)}"

    def test_read_all_notes_structure(self, fixture_path: Path) -> str:
        prs = Presentation(str(fixture_path))
        all_notes = edit_pptx.read_all_notes(prs)
        assert isinstance(all_notes, list), "read_all_notes did not return a list"
        assert len(all_notes) == len(prs.slides), f"Expected {len(prs.slides)} notes entries, got {len(all_notes)}"
        for index, entry in enumerate(all_notes):
            assert isinstance(entry, dict), f"Entry {index} is not a dict"
            assert entry.get("slide_index") == index, f"Unexpected slide_index for entry {index}: {entry}"
            assert isinstance(entry.get("notes"), str), f"Entry {index} notes is not a string"
        return f"entries={len(all_notes)}"

    def test_notes_idempotency(self, fixture_path: Path) -> str:
        fixture_name = fixture_path.name
        out = _output_path(self.ctx, fixture_name, "notes_idempotency.pptx")
        note_text = "Idempotency test note"
        prs = Presentation(str(fixture_path))
        for slide in prs.slides:
            create_pptx.set_notes(slide, note_text)
        create_pptx.save(prs, out)
        prs2 = Presentation(str(out))
        mismatches = sum(1 for slide in prs2.slides if edit_pptx.get_notes(slide) != note_text)
        assert mismatches == 0, f"Expected all notes to round-trip, got {mismatches} mismatches"
        validation = validate_pptx.validate_pptx(out)
        assert validation.valid, f"Output invalid after notes round-trip: {validation.errors}"
        return f"slides={len(prs2.slides)} mismatches={mismatches}"

    def test_reorder_stability(self, fixture_path: Path) -> str:
        fixture_name = fixture_path.name
        prs = Presentation(str(fixture_path))
        if len(prs.slides) < 2:
            raise TestSkipped("presentation has fewer than two slides")
        original = _slide_titles(prs)
        edit_pptx.reorder_slides(prs, len(prs.slides) - 1, 0)
        edit_pptx.reorder_slides(prs, 0, len(prs.slides) - 1)
        out = _output_path(self.ctx, fixture_name, "reorder_stability.pptx")
        prs.save(str(out))
        restored = _slide_titles(Presentation(str(out)))
        assert restored == original, f"Slide order did not restore: {original!r} != {restored!r}"
        return f"slides={len(original)} restored=1"

    def test_multi_image(self, fixture_path: Path) -> str:
        fixture_name = fixture_path.name
        prs = Presentation(str(fixture_path))
        slide_index, slide, pictures = _find_multi_image_slide(prs)
        original_ext = pictures[0].image.ext
        suffix, payload = _replacement_bytes_for_ext(original_ext)
        image_path = _write_binary(_output_path(self.ctx, fixture_name, f"replacement{suffix}"), payload)
        replaced = edit_pptx.replace_image_by_index(slide, 0, image_path)
        assert replaced, "replace_image_by_index returned False"
        out = _output_path(self.ctx, fixture_name, "multi_image_replace.pptx")
        prs.save(str(out))
        prs2 = Presentation(str(out))
        reloaded_pictures = _picture_shapes(prs2.slides[slide_index])
        assert len(reloaded_pictures) == len(pictures), (
            f"Picture count changed on slide {slide_index}: {len(pictures)} -> {len(reloaded_pictures)}"
        )
        validation = validate_pptx.validate_pptx(out)
        assert validation.valid, f"Output invalid after image replacement: {validation.errors}"
        return f"slide={slide_index} pictures={len(pictures)} ext={original_ext}"

    def test_no_placeholder_slide(self, fixture_path: Path) -> str:
        prs = Presentation(str(fixture_path))
        blank_count = sum(
            1
            for slide in prs.slides
            if not any(getattr(shape, "is_placeholder", False) for shape in slide.shapes)
        )
        return f"slides_without_placeholders={blank_count} total_slides={len(prs.slides)}"

    def test_validate(self, fixture_path: Path) -> str:
        result = validate_pptx.validate_pptx(fixture_path)
        assert result.valid, f"validate_pptx failed: {result.errors}"
        return (
            f"slides={result.info.get('slide_count', '?')} "
            f"notes={result.info.get('notes_slide_count', '?')} "
            f"images={result.info.get('image_count', '?')}"
        )

    def test_g1(self) -> str:
        result = vba_runner.check_powerpnt()
        assert isinstance(result, str), "check_powerpnt() did not return a string"
        return f"found: {result}" if result else "not on PATH (acceptable in CI)"

    def test_g2(self) -> str:
        out_path = self.ctx.output_dir / "vba_runner_g2_test.vbs"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out = vba_runner.generate_vbs_runner(
            "dummy.pptx",
            'MsgBox "hello"',
            "BuildSlides",
            out_path,
        )
        content = out.read_text(encoding="utf-8")
        assert "CreateObject" in content, "VBScript missing CreateObject"
        assert "PowerPoint.Application" in content, "VBScript missing PowerPoint.Application"
        assert "BuildSlides" in content, "VBScript missing sub name"
        return f"vbs_lines={len(content.splitlines())}"

    def test_g3(self) -> str:
        out_path = self.ctx.output_dir / "vba_runner_g3_test.bas"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        vba_code = 'Sub Hello()\n    MsgBox "hi"\nEnd Sub'
        out = vba_runner.write_bas_file(vba_code, out_path)
        content = out.read_text(encoding="utf-8")
        assert "Attribute VB_Name" in content, ".bas file missing VB_Name attribute"
        assert "Option Explicit" in content, ".bas file missing Option Explicit"
        assert 'MsgBox "hi"' in content, ".bas file missing VBA body"
        return f"bas_bytes={out.stat().st_size}"


def cleanup(ctx: RunContext) -> None:
    artifacts = [ctx.repo_dir, ctx.output_dir]
    for path in artifacts:
        if path.exists():
            shutil.rmtree(path, ignore_errors=True)
            log_event(ctx, "info", "cleanup.path_removed", path=str(path))
    ctx.cleanup_completed = True
    write_metadata(ctx, status="cleanup-complete")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run the Apache POI-backed pptx regression workflow")
    parser.add_argument("--repo-url", default=DEFAULT_REPO_URL, help="Apache POI git URL")
    parser.add_argument("--ref", default=DEFAULT_REF, help="Git ref to test against")
    parser.add_argument("--keep-artifacts", action="store_true", help="Keep temp clone and outputs for debugging")
    parser.add_argument("--timeout", type=int, default=DEFAULT_TIMEOUT_SECONDS, help="Per-command timeout in seconds")
    parser.add_argument("--retries", type=int, default=DEFAULT_RETRIES, help="Retry count for clone/fetch commands")
    return parser.parse_args()


def build_context(args: argparse.Namespace) -> RunContext:
    TEMP_ROOT.mkdir(parents=True, exist_ok=True)
    run_root = Path(tempfile.mkdtemp(prefix="pptx-poi-suite-", dir=str(TEMP_ROOT)))
    output_dir = run_root / "outputs"
    metadata_path = run_root / "run-metadata.json"
    log_path = run_root / "workflow.log"
    repo_dir = run_root / "apache-poi"
    ctx = RunContext(
        args=args,
        run_root=run_root,
        repo_dir=repo_dir,
        output_dir=output_dir,
        log_path=log_path,
        metadata_path=metadata_path,
        repo_url=args.repo_url,
        ref=args.ref,
    )
    ctx.test_command = " ".join([sys.executable, str(Path(__file__).resolve()), *sys.argv[1:]])
    return ctx


def main() -> int:
    args = parse_args()
    ctx = build_context(args)
    status = "failed"
    error_message: str | None = None
    exit_code = 1
    log_event(ctx, "info", "workflow.start", command=ctx.test_command, run_root=str(ctx.run_root))
    try:
        prepare_poi_baseline(ctx)
        verify_baseline(ctx)
        suite = TestSuite(ctx)
        suite.run()
        failed = [rec for rec in ctx.test_records if rec.status == FAIL]
        status = "passed" if not failed else "failed"
        exit_code = 0 if not failed else 1
        log_event(
            ctx,
            "info" if not failed else "warning",
            "workflow.summary",
            total=len(ctx.test_records),
            passed=sum(1 for rec in ctx.test_records if rec.status == PASS),
            failed=len(failed),
            skipped=sum(1 for rec in ctx.test_records if rec.status == SKIP),
        )
        return exit_code
    except WorkflowError as exc:
        error_message = str(exc)
        log_event(ctx, "error", "workflow.error", error=error_message)
        return 1
    except Exception as exc:
        error_message = f"Unhandled error: {type(exc).__name__}: {exc}"
        log_event(ctx, "error", "workflow.exception", error=error_message, traceback=traceback.format_exc())
        return 1
    finally:
        write_metadata(ctx, status=status, error=error_message)
        if args.keep_artifacts:
            log_event(ctx, "info", "cleanup.skipped", run_root=str(ctx.run_root))
        else:
            cleanup(ctx)


if __name__ == "__main__":
    sys.exit(main())
