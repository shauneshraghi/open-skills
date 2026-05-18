"""Edit existing PowerPoint presentations.

Design decisions from primary sources:
- Slide reorder/delete: python-pptx exposes no API; we manipulate <p:sldIdLst>
  children directly (ECMA-376 §19.2.1.38).  Order of <p:sldId> children defines
  presentation order; no other XML changes are needed for reordering.
- Slide deletion: remove the <p:sldId> element AND drop the OPC relationship from
  the presentation part so the slide XML is no longer referenced.
- Image replacement: update `r:embed` on <a:blip> to a new rId obtained via
  slide.part.get_or_add_image_part() (python-pptx docs §Working with Pictures).
- Notes: slide.notes_slide auto-creates the NotesSlidePart when first accessed
  (python-pptx docs §Working with Notes).  For raw XML inspection we use lxml on
  slide.notes_slide._element (ECMA-376 §19.3.1.43 – notes body placeholder).
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches


# ────────────────────────── slide management ──────────────────────────────

def reorder_slides(prs: Presentation, from_index: int, to_index: int) -> None:
    """Move the slide at `from_index` to `to_index` (0-based).

    Implementation: manipulate <p:sldIdLst> children directly.
    ECMA-376 §19.2.1.38: document order of <p:sldId> elements inside
    <p:sldIdLst> determines presentation slide order.
    """
    n = len(prs.slides)
    if not (0 <= from_index < n and 0 <= to_index < n):
        raise IndexError(f"Slide index out of range (0..{n-1}): from={from_index} to={to_index}")
    if from_index == to_index:
        return

    sldIdLst = prs.part._element.sldIdLst
    sldIds = list(sldIdLst)
    elem = sldIds[from_index]
    sldIdLst.remove(elem)
    insert_pos = to_index if to_index < from_index else to_index
    if insert_pos >= len(list(sldIdLst)):
        sldIdLst.append(elem)
    else:
        sldIdLst.insert(insert_pos, elem)


def remove_slide(prs: Presentation, index: int) -> None:
    """Delete the slide at `index` (0-based) from the presentation.

    Steps:
    1. Find and remove the <p:sldId> element from <p:sldIdLst>.
    2. Drop the OPC relationship from the presentation part so the slide
       XML part is no longer referenced (python-pptx OPC model).

    Note: any notes slide associated with the removed slide is *not*
    automatically cleaned up here; a save/reload cycle will compact the package.
    """
    n = len(prs.slides)
    if not (0 <= index < n):
        raise IndexError(f"Slide index {index} out of range (0..{n-1})")

    sldIdLst = prs.part._element.sldIdLst
    sldIds = list(sldIdLst)
    target = sldIds[index]
    rId = target.get(qn("r:id"))
    sldIdLst.remove(target)
    prs.part.drop_rel(rId)


def get_slide_count(prs: Presentation) -> int:
    """Return the number of slides in the presentation."""
    return len(prs.slides)


# ────────────────────────── notes operations ──────────────────────────────

def set_notes(slide: Any, text: str) -> None:
    """Set (replace) the speaker notes text on a slide.

    python-pptx auto-creates the NotesSlidePart (and a notesMaster if absent)
    when slide.notes_slide is first accessed.  The text body placeholder has
    ph type='body' idx=1 per ECMA-376 §19.3.1.43.
    """
    ns = slide.notes_slide
    ns.notes_text_frame.text = text


def get_notes(slide: Any) -> str:
    """Return the plain-text speaker notes for a slide, or '' if none."""
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


def set_notes_paragraphs(slide: Any, paragraphs: list[str]) -> None:
    """Set notes as multiple paragraphs using direct lxml manipulation.

    Use this instead of set_notes() when you need to preserve empty lines,
    set per-run formatting, or repair malformed notes XML from third-party tools.

    ECMA-376 §20.1.2.1.24: each <a:p> is a paragraph; <a:r> is a run;
    <a:t> carries the text content.
    """
    ns = slide.notes_slide
    elm = ns._element
    body_sp = _find_notes_body_sp(elm)
    if body_sp is None:
        ns.notes_text_frame.text = "\n".join(paragraphs)
        return

    txBody = body_sp.find(qn("p:txBody"))
    if txBody is None:
        return

    for p_elem in txBody.findall(qn("a:p")):
        txBody.remove(p_elem)

    for text in paragraphs:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        if text:
            r_elem = etree.SubElement(p_elem, qn("a:r"))
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = text
        else:
            etree.SubElement(p_elem, qn("a:endParaRPr"))


def _find_notes_body_sp(notes_elm: Any) -> Any | None:
    """Find the <p:sp> with ph type='body' in the notes shape tree.

    The notes slide has two required placeholders (ECMA-376 §19.3.1.43):
    - sldImg (idx=0): slide thumbnail
    - body   (idx=1): notes text
    """
    cSld = notes_elm.find(qn("p:cSld"))
    if cSld is None:
        return None
    spTree = cSld.find(qn("p:spTree"))
    if spTree is None:
        return None
    for sp in spTree.findall(qn("p:sp")):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        ph = nvPr.find(qn("p:ph"))
        if ph is not None and ph.get("type") == "body":
            return sp
    return None


def read_all_notes(prs: Presentation) -> list[dict]:
    """Return a list of dicts {'slide_index': int, 'notes': str} for all slides."""
    result = []
    for i, slide in enumerate(prs.slides):
        result.append({"slide_index": i, "notes": get_notes(slide)})
    return result


# ────────────────────────── image operations ──────────────────────────────

def replace_image(slide: Any, shape_name: str, new_image_path: str | Path) -> bool:
    """Replace the image in a named Picture shape by swapping its OPC relationship.

    Approach (python-pptx + lxml):
    1. Find the shape by name.
    2. Locate the <a:blip r:embed="rId…"> inside <p:blipFill>.
    3. Add the new image as a new part via slide.part.get_or_add_image_part().
    4. Update r:embed on <a:blip> to point to the new rId.

    The old image part remains in the package until save/reload; this is
    harmless and conforming per ECMA-376 Part 2 §11.3.
    """
    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        if shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE
            return False
        blip = shape._element.blipFill.blip
        _, new_rId = slide.part.get_or_add_image_part(str(new_image_path))
        blip.set(qn("r:embed"), new_rId)
        return True
    return False


def replace_image_by_index(slide: Any, picture_index: int, new_image_path: str | Path) -> bool:
    """Replace the nth picture (0-based) on a slide."""
    pics = [s for s in slide.shapes if s.shape_type == 13]
    if not (0 <= picture_index < len(pics)):
        return False
    blip = pics[picture_index]._element.blipFill.blip
    _, new_rId = slide.part.get_or_add_image_part(str(new_image_path))
    blip.set(qn("r:embed"), new_rId)
    return True


def set_alt_text(shape: Any, alt_text: str) -> None:
    """Set accessibility alt text on a shape.

    ECMA-376 §20.1.2.2.8: <p:cNvPr descr="…"> or <a:cNvPr descr="…"> carries
    the accessibility description.  No python-pptx high-level API exists for this.
    """
    for attr in ("nvPicPr", "nvSpPr", "nvGraphicFramePr"):
        nvPr_elem = getattr(shape._element, attr, None)
        if nvPr_elem is not None:
            cNvPr = getattr(nvPr_elem, "cNvPr", None)
            if cNvPr is not None:
                cNvPr.set("descr", alt_text)
                return


def get_alt_text(shape: Any) -> str:
    """Return the alt text (descr) from a shape's cNvPr element, or ''."""
    for attr in ("nvPicPr", "nvSpPr", "nvGraphicFramePr"):
        nvPr_elem = getattr(shape._element, attr, None)
        if nvPr_elem is not None:
            cNvPr = getattr(nvPr_elem, "cNvPr", None)
            if cNvPr is not None:
                return cNvPr.get("descr", "")
    return ""


# ────────────────────────── text operations ───────────────────────────────

def update_placeholder_text(slide: Any, placeholder_idx: int, text: str) -> bool:
    """Replace text in a placeholder by index."""
    try:
        slide.placeholders[placeholder_idx].text = text
        return True
    except (KeyError, IndexError):
        return False


def find_shapes_by_type(slide: Any, shape_type: int) -> list:
    """Return all shapes on a slide matching shape_type (use MSO_SHAPE_TYPE ints)."""
    return [s for s in slide.shapes if s.shape_type == shape_type]


# ────────────────────────── CLI ──────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="Edit an existing .pptx file")
    sub = parser.add_subparsers(dest="cmd", required=True)

    p_reorder = sub.add_parser("reorder", help="Move a slide to a new position")
    p_reorder.add_argument("file", type=Path)
    p_reorder.add_argument("from_index", type=int)
    p_reorder.add_argument("to_index", type=int)
    p_reorder.add_argument("--out", type=Path, default=None)

    p_remove = sub.add_parser("remove", help="Delete a slide")
    p_remove.add_argument("file", type=Path)
    p_remove.add_argument("index", type=int)
    p_remove.add_argument("--out", type=Path, default=None)

    p_notes = sub.add_parser("set-notes", help="Set speaker notes on a slide")
    p_notes.add_argument("file", type=Path)
    p_notes.add_argument("slide_index", type=int)
    p_notes.add_argument("text", type=str)
    p_notes.add_argument("--out", type=Path, default=None)

    p_read = sub.add_parser("read-notes", help="Print all speaker notes")
    p_read.add_argument("file", type=Path)

    p_img = sub.add_parser("replace-image", help="Replace a named picture on a slide")
    p_img.add_argument("file", type=Path)
    p_img.add_argument("slide_index", type=int)
    p_img.add_argument("shape_name", type=str)
    p_img.add_argument("new_image", type=Path)
    p_img.add_argument("--out", type=Path, default=None)

    args = parser.parse_args()
    prs = Presentation(str(args.file))
    out = args.out or args.file

    if args.cmd == "reorder":
        reorder_slides(prs, args.from_index, args.to_index)
        prs.save(str(out))
        print(f"Reordered: slide {args.from_index} → position {args.to_index}. Saved: {out}")

    elif args.cmd == "remove":
        remove_slide(prs, args.index)
        prs.save(str(out))
        print(f"Removed slide {args.index}. Saved: {out}")

    elif args.cmd == "set-notes":
        slide = prs.slides[args.slide_index]
        set_notes(slide, args.text)
        prs.save(str(out))
        print(f"Set notes on slide {args.slide_index}. Saved: {out}")

    elif args.cmd == "read-notes":
        for entry in read_all_notes(prs):
            idx = entry["slide_index"]
            notes = entry["notes"]
            if notes:
                print(f"Slide {idx}: {notes!r}")
            else:
                print(f"Slide {idx}: (no notes)")

    elif args.cmd == "replace-image":
        slide = prs.slides[args.slide_index]
        ok = replace_image(slide, args.shape_name, args.new_image)
        if ok:
            prs.save(str(out))
            print(f"Image replaced. Saved: {out}")
        else:
            print(f"Shape '{args.shape_name}' not found on slide {args.slide_index}.")


if __name__ == "__main__":
    main()
